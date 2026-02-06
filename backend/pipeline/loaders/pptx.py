"""PowerPoint document loader using python-pptx."""

import asyncio
from pathlib import Path

from pptx import Presentation

from .base import BaseLoader, LoadedDocument


class PPTXLoader(BaseLoader):
    """Load PPTX files using python-pptx."""

    supported_extensions = [".pptx"]

    async def load(self, file_path: str | Path) -> LoadedDocument:
        path = self._validate_file(file_path)
        return await asyncio.to_thread(self._load_sync, path)

    def _load_sync(self, path: Path) -> LoadedDocument:
        prs = Presentation(str(path))
        pages = []
        all_text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    table_text = self._extract_table(shape.table)
                    if table_text:
                        slide_texts.append(table_text)

            if slide_texts:
                slide_content = "\n".join(slide_texts)
                pages.append({
                    "page_num": slide_num,
                    "content": slide_content,
                })
                all_text_parts.append(f"[슬라이드 {slide_num}]\n{slide_content}")

        return LoadedDocument(
            content="\n\n".join(all_text_parts),
            metadata={
                "filename": path.name,
                "file_type": "pptx",
                "slide_count": len(prs.slides),
            },
            pages=pages,
        )

    def _extract_table(self, table) -> str:
        """Extract table from slide."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
